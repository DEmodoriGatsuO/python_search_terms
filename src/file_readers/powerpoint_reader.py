import logging
from pptx import Presentation

def read_powerpoint(file_path):
    """
    PowerPointファイルからテキストを抽出する関数

    Args:
        file_path (str): PowerPointファイルのパス

    Returns:
        str: 抽出されたテキスト
    """
    try:
        presentation = Presentation(file_path)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logging.error(f"Error reading PowerPoint file {file_path}: {e}")
        return ""